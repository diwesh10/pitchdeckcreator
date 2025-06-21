import streamlit as st
import google.generativeai as genai
import os
from pptx import Presentation
from pptx.util import Inches
import io

# --- Configuration ---
st.set_page_config(
    page_title="AutoPitch Pro",
    page_icon="🚀",
    layout="centered"
)

# --- App Title ---
st.title("AutoPitch Pro – AI-Powered Startup Pitch Generator")
st.header("Turn your one-line startup idea into a complete pitch deck.")

# --- API Key Management ---
try:
    # Get API key from Streamlit secrets
    api_key = st.secrets["GOOGLE_API_KEY"]
    genai.configure(api_key=api_key)
except (KeyError, FileNotFoundError):
    st.error("🛑 **API key not found!** Please add your Google API Key to your Streamlit secrets.")
    st.info("For local development, create a `.streamlit/secrets.toml` file and add your key: `GOOGLE_API_KEY = 'YOUR_KEY_HERE'`")
    st.stop()


# --- Main App ---
startup_idea = st.text_input("Enter your startup idea (e.g., 'An app that uses AI to create personalized workout plans')")

def generate_pitch_deck(idea):
    prompt = f"""
    You are an expert startup consultant. Create a 10-slide pitch deck for the following startup idea: "{idea}"

    Please provide the content for each of the following slides, using "---SLIDE---" as a separator between each slide.
    For each slide, the first line should be the slide title, and the following lines should be bullet points starting with a hyphen '-'.

    The slides are:
    1. Title Slide: Company Name and Tagline
    2. Problem
    3. Solution
    4. Market Size
    5. Product
    6. Business Model
    7. Go-to-Market Strategy
    8. Competition
    9. Team
    10. Call to Action / Ask
    """
    model = genai.GenerativeModel('gemini-1.5-flash-latest')
    response = model.generate_content(prompt)
    return response.text

def create_pptx(pitch_content):
    prs = Presentation()
    slides = pitch_content.strip().split("---SLIDE---")

    for i, slide_text in enumerate(slides):
        if not slide_text.strip():
            continue

        lines = slide_text.strip().split('\n')
        title = lines[0].strip()
        content_points = lines[1:]

        if i == 0: # Title Slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title_placeholder = slide.shapes.title
            subtitle_placeholder = slide.placeholders[1]
            title_placeholder.text = title
            subtitle_placeholder.text = "\n".join(point.strip("- ").strip() for point in content_points)
        else: # Content Slide
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_placeholder = slide.shapes.title
            body_placeholder = slide.placeholders[1]
            title_placeholder.text = title
            tf = body_placeholder.text_frame
            tf.clear() # Clear existing text
            for point in content_points:
                p = tf.add_paragraph()
                p.text = point.strip("- ").strip()
                p.level = 0


    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io


if st.button("Generate Pitch"):
    if not startup_idea:
        st.error("Please enter your startup idea.")
    else:
        try:
            with st.spinner("Generating your pitch deck... this may take a moment..."):
                pitch_content = generate_pitch_deck(startup_idea)
                pptx_file = create_pptx(pitch_content)

            st.subheader("Your Pitch Deck is Ready!")
            st.markdown("Click the button below to download your presentation.")

            st.download_button(
                label="Download Pitch Deck (.pptx)",
                data=pptx_file,
                file_name=f"{startup_idea[:20].replace(' ','_')}_pitch_deck.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )

            st.markdown("---")
            st.subheader("Generated Raw Text:")
            st.markdown(pitch_content)

        except Exception as e:
            st.error(f"An error occurred: {e}")
